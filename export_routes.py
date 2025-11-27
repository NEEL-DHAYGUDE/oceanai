from flask import Blueprint, send_file, jsonify
from models import Section, Project
from docx import Document
from pptx import Presentation
import os
import re

export_bp = Blueprint("export", __name__)

# Utility: sanitize filename (important!)
def safe_filename(name):
    return re.sub(r'[\\/*?:"<>|]', "_", name)

@export_bp.get("/export/docx/<project_id>")
def export_docx(project_id):
    # --- VALIDATION FIX ---
    if not project_id or project_id == "undefined":
        return jsonify({"error": "Invalid Project ID"}), 400

    project = Project.query.get(project_id)
    if not project:
        return jsonify({"error": "Project not found"}), 404
    # ----------------------

    sections = Section.query.filter_by(project_id=project_id).all()

    doc = Document()
    doc.add_heading(project.title, level=1)

    for sec in sections:
        doc.add_heading(sec.title, level=2)
        doc.add_paragraph(sec.content if sec.content else "")

    # filename fix
    filename = safe_filename(f"{project.title}.docx")
    filepath = os.path.join("generated", filename)
    os.makedirs("generated", exist_ok=True)

    doc.save(filepath)
    return send_file(filepath, as_attachment=True)


@export_bp.get("/export/pptx/<project_id>")
def export_pptx(project_id):
    # --- VALIDATION FIX ---
    if not project_id or project_id == "undefined":
        return jsonify({"error": "Invalid Project ID"}), 400

    project = Project.query.get(project_id)
    if not project:
        return jsonify({"error": "Project not found"}), 404
    # ----------------------

    sections = Section.query.filter_by(project_id=project_id).all()

    pres = Presentation()

    # Remove default first slide (fix double slide issue)
    if len(pres.slides) > 0:
        slide = pres.slides[0]
        pres.slides._sldIdLst.remove(slide._element.getparent())

    for sec in sections:
        slide = pres.slides.add_slide(pres.slide_layouts[1])
        slide.shapes.title.text = sec.title
        slide.placeholders[1].text = sec.content if sec.content else ""

    filename = safe_filename(f"{project.title}.pptx")
    filepath = os.path.join("generated", filename)
    os.makedirs("generated", exist_ok=True)

    pres.save(filepath)
    return send_file(filepath, as_attachment=True)
